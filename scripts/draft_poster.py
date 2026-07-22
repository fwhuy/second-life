from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "poster.pptx"
OUTPUT = ROOT / "poster_draft.pptx"
REPORTS = ROOT / "model" / "reports" / "baseline_resnet50"

PURPLE = RGBColor(91, 33, 145)
GREEN = RGBColor(47, 90, 60)
SAGE = RGBColor(220, 231, 211)
ORANGE = RGBColor(217, 123, 41)
BLUE = RGBColor(43, 108, 176)
INK = RGBColor(35, 55, 43)
MUTED = RGBColor(92, 107, 96)
PAPER = RGBColor(255, 253, 246)
BG = RGBColor(242, 240, 246)
LINE = RGBColor(218, 214, 224)
WHITE = RGBColor(255, 255, 255)


def add_box(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False,
             font="Arial", align=PP_ALIGN.LEFT, margin=0.08,
             valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def section_title(slide, number, title, x, y, w):
    add_text(slide, number, x, y, .7, .55, 24, ORANGE, True, "Arial")
    add_text(slide, title, x + .65, y - .02, w - .65, .62, 25, PURPLE, True, "Arial")


def metric_card(slide, value, label, x, y, w, color):
    add_box(slide, x, y, w, 1.55, PAPER, RGBColor(229, 220, 200))
    add_text(slide, value, x + .12, y + .10, w - .24, .72, 36, color, True, "Georgia", PP_ALIGN.CENTER)
    add_text(slide, label, x + .12, y + .84, w - .24, .48, 13, MUTED, True, "Arial", PP_ALIGN.CENTER)


def pipeline_step(slide, num, title, sub, x, y, w=3.05):
    add_box(slide, x, y, w, 2.0, PAPER, RGBColor(201, 189, 159))
    add_text(slide, num, x + .15, y + .12, .55, .55, 20, WHITE, True, "Arial", PP_ALIGN.CENTER, 0,
             MSO_ANCHOR.MIDDLE).fill.solid()
    badge = slide.shapes[-1]
    badge.fill.fore_color.rgb = GREEN
    add_text(slide, title, x + .78, y + .12, w - .9, .50, 18, INK, True)
    add_text(slide, sub, x + .18, y + .78, w - .36, 1.0, 13, MUTED)


def main():
    prs = Presentation(SOURCE)
    slide = prs.slides[0]
    # Preserve the supplied canvas/theme and embedded fonts, but replace all
    # example-project objects with the trash-classification poster draft.
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)

    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = BG

    # Header
    add_box(slide, .55, .45, 40.20, 2.75, WHITE, WHITE, False)
    add_text(slide, "SECOND LIFE AI", .95, .69, 8.0, .48, 17, ORANGE, True, "Arial")
    add_text(slide, "Leakage-Aware Trash Classification", .92, 1.08, 28.8, 1.0, 40, INK, True, "Georgia")
    add_text(slide, "A reproducible six-class pipeline, an honest held-out test, and a live sorting experience",
             .96, 2.13, 29.0, .48, 17, MUTED, False)
    add_text(slide, "GROUP [##]  ·  [NAME]  ·  [NAME]  ·  [NAME]  ·  [NAME]",
             29.65, 1.02, 10.2, .85, 17, PURPLE, True, "Arial", PP_ALIGN.RIGHT)
    add_text(slide, "NYU SHANGHAI · AI SUMMER PROGRAM 2026",
             29.70, 2.04, 10.1, .42, 13, MUTED, True, "Arial", PP_ALIGN.RIGHT)

    # Main panels
    add_box(slide, .55, 3.55, 10.35, 26.65, WHITE, LINE, False)
    add_box(slide, 11.25, 3.55, 18.05, 26.65, WHITE, LINE, False)
    add_box(slide, 29.65, 3.55, 11.10, 17.20, WHITE, LINE, False)
    add_box(slide, 29.65, 21.10, 11.10, 9.10, WHITE, LINE, False)

    # 1 Problem & Data
    section_title(slide, "01", "Problem & Data", .86, 3.88, 9.6)
    add_text(slide, "Can a model sort a waste photo reliably—without seeing the same physical object in both training and evaluation?",
             .92, 4.62, 9.55, 1.65, 20, INK, True, "Georgia")
    add_text(slide, "TrashNet · 2,527 images · 6 classes", .92, 6.35, 9.2, .55, 17, GREEN, True)

    counts = [("paper", 594), ("glass", 501), ("plastic", 482), ("metal", 410), ("cardboard", 403), ("trash", 137)]
    max_count = max(v for _, v in counts)
    y = 7.12
    for label, value in counts:
        add_text(slide, label.title(), .96, y, 2.05, .38, 13, MUTED, True)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.0), Inches(y + .04),
                                     Inches(5.6 * value / max_count), Inches(.26))
        bar.fill.solid(); bar.fill.fore_color.rgb = GREEN if label != "trash" else ORANGE
        bar.line.fill.background()
        add_text(slide, str(value), 8.72, y - .01, .85, .36, 13, INK, True, align=PP_ALIGN.RIGHT)
        y += .62

    add_box(slide, .92, 11.12, 9.55, 2.35, RGBColor(246, 232, 216), RGBColor(233, 196, 138))
    add_text(slide, "THE DATASET TRAP", 1.15, 11.34, 3.6, .42, 14, ORANGE, True)
    add_text(slide, "TrashNet contains many photographs of the same object from different angles. Random image-level splits can leak object identity across sets.",
             1.15, 11.82, 8.95, 1.22, 16, INK)

    # 2 Evaluation
    section_title(slide, "02", "Evaluation Contract", .86, 14.03, 9.6)
    add_text(slide, "Split by near-duplicate group, augment only after splitting, and touch the quarantined test set once.",
             .96, 14.78, 9.45, 1.42, 18, INK, True, "Georgia")
    metric_card(slide, "1,732", "TRAIN", .96, 16.45, 2.75, GREEN)
    metric_card(slide, "434", "VALIDATION", 3.88, 16.45, 2.75, BLUE)
    metric_card(slide, "361", "HELD-OUT TEST", 6.80, 16.45, 3.25, ORANGE)
    add_text(slide, "Test measured once before intervention; never used to select experiments.",
             1.04, 18.25, 8.95, .82, 14, MUTED)

    add_box(slide, .92, 19.35, 9.55, 4.28, SAGE, RGBColor(175, 198, 173))
    add_text(slide, "WHY THIS IS CREATIVE", 1.16, 19.62, 4.6, .42, 14, GREEN, True)
    creative = [
        "Accuracy that survives scrutiny",
        "Fixed seeds and committed split identities",
        "Every run logged with a configuration hash",
        "Test-spent guard blocks accidental re-evaluation",
    ]
    for i, line in enumerate(creative):
        add_text(slide, "•", 1.18, 20.12 + i * .58, .34, .40, 17, PURPLE, True)
        add_text(slide, line, 1.52, 20.10 + i * .58, 8.34, .44, 15, INK, i == 0)

    add_text(slide, "OFFICIAL METRIC", 1.02, 24.08, 3.4, .38, 13, ORANGE, True)
    add_text(slide, "89.20%", 1.00, 24.43, 5.0, 1.10, 42, PURPLE, True, "Georgia")
    add_text(slide, "322 / 361 held-out test images", 1.03, 25.51, 8.8, .55, 17, INK, True)
    add_text(slide, "This is the test result. The higher 94.93% figure is validation + TTA.",
             1.03, 26.15, 8.95, 1.0, 14, MUTED)

    # 3 Method
    section_title(slide, "03", "Leakage-Free Method", 11.62, 3.88, 17.1)
    add_text(slide, "A simple model with a strict evaluation pipeline", 11.72, 4.58, 16.7, .58, 20, GREEN, True, "Georgia")
    steps = [
        ("1", "Deduplicate", "Perceptual groups keep near-identical views together."),
        ("2", "Split", "Train / validation / test are group-disjoint."),
        ("3", "Augment", "Basic transforms apply to training images only."),
        ("4", "Fine-tune", "ImageNet ResNet-50, 224px, two-phase training."),
        ("5", "Infer", "Average four deterministic views with TTA."),
    ]
    xs = [11.72, 15.08, 18.44, 21.80, 25.16]
    for (num, title, sub), x in zip(steps, xs):
        pipeline_step(slide, num, title, sub, x, 5.45)
    for x in [14.79, 18.15, 21.51, 24.87]:
        add_text(slide, "→", x, 6.04, .34, .5, 19, ORANGE, True, align=PP_ALIGN.CENTER)

    add_box(slide, 11.72, 7.84, 16.85, 3.05, RGBColor(30, 58, 43), RGBColor(30, 58, 43))
    add_text(slide, "MODEL", 12.05, 8.14, 2.0, .4, 13, RGBColor(157, 184, 160), True)
    add_text(slide, "ResNet-50 · 224 px", 12.02, 8.56, 5.0, .65, 25, WHITE, True, "Georgia")
    add_text(slide, "weighted sampler  ·  label smoothing 0.1  ·  weight decay 0.05  ·  seed 42",
             12.04, 9.38, 15.75, .52, 15, RGBColor(217, 228, 214))
    add_text(slide, "The model is intentionally conventional; the contribution is trustworthy measurement.",
             12.04, 10.05, 15.75, .45, 14, RGBColor(240, 178, 122), True)

    section_title(slide, "", "What changed the result?", 11.70, 11.35, 16.8)
    # Compact experiment table
    headers = ["Configuration", "Val. accuracy", "Net images", "Decision"]
    rows = [
        ["Raw baseline", "93.09%", "—", "Reference"],
        ["+ four-view TTA", "94.93%", "+8", "Adopted"],
        ["384px progressive", "94.93%", "0", "Rejected"],
        ["Modern augmentation", "95.62%", "+3 vs TTA", "Rejected"],
        ["Weight decay 0.10", "94.70%", "−1", "Rejected"],
    ]
    x0, y0 = 11.80, 12.18
    widths = [6.0, 3.4, 3.0, 3.7]
    x = x0
    for h, w in zip(headers, widths):
        add_box(slide, x, y0, w, .62, PURPLE, PURPLE, False)
        add_text(slide, h, x + .08, y0 + .10, w - .16, .36, 13, WHITE, True, align=PP_ALIGN.CENTER)
        x += w
    for ri, row in enumerate(rows):
        x = x0; yy = y0 + .68 + ri * .67
        for ci, (cell, w) in enumerate(zip(row, widths)):
            fill = RGBColor(242, 247, 239) if ri == 1 else PAPER
            add_box(slide, x, yy, w, .62, fill, LINE, False)
            color = GREEN if ri == 1 else INK
            add_text(slide, cell, x + .08, yy + .10, w - .16, .35, 13, color, ri == 1,
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            x += w
    add_text(slide, "Pre-registered acceptance gate: at least +7 validation images, no macro-F1 drop, and no major class-recall loss.",
             11.88, 16.38, 16.4, .58, 13, MUTED)

    # Visual evidence
    section_title(slide, "", "Training behavior & error evidence", 11.70, 17.20, 16.8)
    curves = REPORTS / "curves_fold0.png"
    errors = REPORTS / "step1_diagnosis" / "top30_highest_loss_1.png"
    if curves.exists():
        slide.shapes.add_picture(str(curves), Inches(11.82), Inches(18.05), width=Inches(8.05), height=Inches(6.3))
    if errors.exists():
        slide.shapes.add_picture(str(errors), Inches(20.16), Inches(18.05), width=Inches(8.05), height=Inches(6.3))
    add_text(slide, "Learning curves", 11.90, 24.46, 7.8, .42, 14, MUTED, True, align=PP_ALIGN.CENTER)
    add_text(slide, "Highest-loss audit examples", 20.24, 24.46, 7.8, .42, 14, MUTED, True, align=PP_ALIGN.CENTER)

    add_box(slide, 11.82, 25.12, 16.38, 3.95, RGBColor(246, 232, 216), RGBColor(233, 196, 138))
    add_text(slide, "ERROR AUDIT", 12.08, 25.40, 3.0, .4, 14, ORANGE, True)
    add_text(slide, "29 of the 30 highest-loss test examples were errors. Manual review judged 14 clearly ambiguous or mislabeled, plus one borderline case.",
             12.08, 25.88, 15.85, 1.22, 18, INK, True, "Georgia")
    add_text(slide, "Interpretation: some remaining error is a data-quality ceiling—not automatically a modelling failure.",
             12.08, 27.42, 15.85, .72, 15, MUTED)

    # 4 Results
    section_title(slide, "04", "Results", 30.02, 3.88, 10.2)
    metric_card(slide, "89.20%", "HELD-OUT TEST · 322/361", 30.05, 4.72, 4.90, ORANGE)
    metric_card(slide, "94.93%", "VALIDATION + 4-VIEW TTA", 35.15, 4.72, 4.90, GREEN)
    metric_card(slide, "94.14%", "VALIDATION MACRO-F1", 30.05, 6.48, 4.90, PURPLE)
    metric_card(slide, "+8", "VAL IMAGES FIXED BY TTA", 35.15, 6.48, 4.90, BLUE)
    confusion = REPORTS / "confusion_baseline_resnet50_val_fold0_tta.png"
    if confusion.exists():
        slide.shapes.add_picture(str(confusion), Inches(30.18), Inches(8.48), width=Inches(9.82), height=Inches(8.20))
    add_text(slide, "Validation confusion matrix · same locked checkpoint + TTA",
             30.18, 16.78, 9.82, .52, 14, MUTED, True, align=PP_ALIGN.CENTER)
    add_box(slide, 30.18, 17.50, 9.82, 2.50, SAGE, RGBColor(175, 198, 173))
    add_text(slide, "KEY RESULT", 30.46, 17.76, 2.4, .38, 13, GREEN, True)
    add_text(slide, "TTA improved accuracy without retraining. All three training interventions failed the pre-registered adoption gate.",
             30.46, 18.18, 9.25, 1.35, 17, INK, True, "Georgia")

    # 5 Takeaways
    section_title(slide, "05", "Takeaways", 30.02, 21.42, 10.2)
    takeaways = [
        "Split by object, not just image.",
        "Report test and validation separately.",
        "Better evaluation can matter more than a newer backbone.",
        "Remaining errors include ambiguity and label noise.",
    ]
    for i, line in enumerate(takeaways):
        yy = 22.22 + i * 1.02
        badge = add_box(slide, 30.18, yy, .62, .62, PURPLE, PURPLE)
        add_text(slide, str(i + 1), 30.18, yy + .04, .62, .42, 15, WHITE, True,
                 align=PP_ALIGN.CENTER, margin=0, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, line, 31.03, yy - .01, 8.60, .72, 17, INK, True, "Georgia")
    add_box(slide, 30.10, 26.74, 9.92, 2.48, RGBColor(30, 58, 43), RGBColor(30, 58, 43))
    add_text(slide, "LIVE DEMO", 30.40, 27.04, 2.1, .38, 13, RGBColor(157, 184, 160), True)
    add_text(slide, "Insight Lab → identify → choose a future", 30.38, 27.48, 7.15, .72, 18, WHITE, True, "Georgia")
    qr = add_box(slide, 37.85, 27.00, 1.80, 1.80, WHITE, WHITE, False)
    add_text(slide, "QR", 37.85, 27.42, 1.80, .45, 20, PURPLE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "[ADD DEMO LINK]", 36.90, 28.84, 3.7, .34, 10, MUTED, True, align=PP_ALIGN.CENTER)

    add_text(slide, "Source: TrashNet (Thung & Yang, 2016). Test result was measured once and not reused for model selection. All values trace to experiments.csv and FINAL_SUMMARY.md.",
             .80, 30.37, 39.70, .34, 10, MUTED, False, "Arial", PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
